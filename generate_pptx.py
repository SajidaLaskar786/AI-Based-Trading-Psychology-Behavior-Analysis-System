"""
TradePsych AI — Pitch Deck Generator
Generates a fully styled 12-slide .pptx file
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colours ──────────────────────────────────────────────
BG          = RGBColor(0x08, 0x10, 0x1f)   # dark navy
BG2         = RGBColor(0x0d, 0x1a, 0x2e)
PRIMARY     = RGBColor(0x7C, 0x5C, 0xFC)   # purple
PRIMARY_L   = RGBColor(0xa3, 0x8b, 0xfd)
ACCENT      = RGBColor(0x22, 0xd3, 0xee)   # cyan
GREEN       = RGBColor(0x10, 0xb9, 0x81)
RED         = RGBColor(0xf4, 0x3f, 0x5e)
YELLOW      = RGBColor(0xf5, 0x9e, 0x0b)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED  = RGBColor(0x94, 0xa3, 0xb8)
TEXT_DIM    = RGBColor(0x64, 0x74, 0x8b)
CARD_BG     = RGBColor(0x12, 0x20, 0x38)

# ── Slide dimensions (16:9 widescreen) ───────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill_color=None, fill_alpha=None,
             line_color=None, line_width=Pt(0), radius=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_multiline(slide, lines, x, y, w, h,
                  default_size=12, default_color=WHITE,
                  default_bold=False, align=PP_ALIGN.LEFT,
                  line_spacing=1.0):
    """lines = list of (text, size, color, bold, italic)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, color, bold, italic = item, default_size, default_color, default_bold, False
        else:
            text  = item[0]
            size  = item[1] if len(item) > 1 else default_size
            color = item[2] if len(item) > 2 else default_color
            bold  = item[3] if len(item) > 3 else default_bold
            italic= item[4] if len(item) > 4 else False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb


def slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_tag(slide, text, x, y):
    """Pill-shaped label"""
    r = add_rect(slide, x, y, len(text)*0.085+0.3, 0.28,
                 fill_color=RGBColor(0x1a, 0x0f, 0x40),
                 line_color=PRIMARY, line_width=Pt(0.5))
    add_text(slide, text, x+0.05, y+0.02, len(text)*0.085+0.2, 0.24,
             font_size=9, bold=True, color=PRIMARY_L,
             align=PP_ALIGN.CENTER)
    return r


def add_logo_bar(slide):
    add_text(slide, "ψ  TradePsych AI", 0.35, 0.18, 2.5, 0.32,
             font_size=11, bold=True, color=TEXT_MUTED)


def add_card(slide, x, y, w, h, fill=CARD_BG, border=None):
    border = border or RGBColor(0x1e, 0x32, 0x50)
    r = add_rect(slide, x, y, w, h, fill_color=fill,
                 line_color=border, line_width=Pt(0.75))
    return r


def add_slide_title(slide, line1, accent_word="", line2=""):
    """Big slide headline with optional accent colour word"""
    y = 0.95
    if accent_word:
        # two-part title
        add_multiline(slide, [
            (line1, 20, WHITE, True),
            (accent_word, 20, PRIMARY_L, True),
            (line2, 20, WHITE, True) if line2 else ("", 1, WHITE, False),
        ], 0.35, y, 12.6, 0.9)
    else:
        add_text(slide, line1, 0.35, y, 12.6, 0.9,
                 font_size=20, bold=True, color=WHITE)


def h_divider(slide, y=1.82):
    r = add_rect(slide, 0.35, y, 0.6, 0.04, fill_color=PRIMARY)
    return r


def check_yes():
    return "✓"

def check_no():
    return "✗"


# ══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)

# Gradient overlay strip
add_rect(s, 0, 0, 13.33, 7.5, fill_color=RGBColor(0x0a, 0x10, 0x25))

# Left accent bar
add_rect(s, 0, 0, 0.08, 7.5, fill_color=PRIMARY)

# Brand name
add_text(s, "TradePsych", 0.6, 1.4, 7, 1.3,
         font_size=60, bold=True, color=WHITE, font_name="Calibri")
add_text(s, "AI", 0.6, 2.65, 3, 0.9,
         font_size=60, bold=True, color=PRIMARY_L, font_name="Calibri")

# Tagline
add_text(s, '"Know Why You Lose Before You Lose Again."',
         0.6, 3.6, 8, 0.5, font_size=16, italic=True, color=TEXT_MUTED)

# Divider
add_rect(s, 0.6, 4.22, 0.8, 0.05, fill_color=ACCENT)

# Sub-copy
add_text(s, "AI-powered trading psychology analysis that turns behavioral\n"
            "blind spots into competitive edge.",
         0.6, 4.4, 8, 0.8, font_size=13, color=TEXT_MUTED)

# Right panel — stats cards
for i, (val, lbl, col) in enumerate([
    ("MVP Live", "Working Product", GREEN),
    ("2025", "Founded", PRIMARY_L),
    ("19+", "Behavioral Metrics", ACCENT),
    ("< 5s", "Report Generated", YELLOW),
]):
    cx = 9.5 + (i % 2) * 1.75
    cy = 2.3 + (i // 2) * 1.4
    add_card(s, cx, cy, 1.6, 1.2, border=col)
    add_text(s, val, cx+0.1, cy+0.15, 1.4, 0.55,
             font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, cx+0.05, cy+0.72, 1.5, 0.35,
             font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Badge strip
add_text(s, "Investor Pitch  ·  2026  ·  Confidential",
         0.6, 6.8, 8, 0.4, font_size=10, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "⚠  THE PROBLEM", 0.35, 0.55)
add_text(s, "80% of retail traders lose money — not because of bad strategy,\nbut because of bad psychology.",
         0.35, 0.9, 12.6, 0.9, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.85)

# Pain cards
pains = [
    ("😤", "Emotional Trading",
     "Traders revenge-trade, panic-exit, and FOMO into bad positions under stress — driven by emotion, not plan."),
    ("📉", "No Self-Awareness",
     "Brokers provide P&L reports but zero behavioral insight. Traders know their losses — not their WHY."),
    ("🔁", "Repeating Mistakes",
     "Without diagnosis, the same patterns repeat across hundreds of trades, costing real money every session."),
]
for i, (icon, title, desc) in enumerate(pains):
    cx = 0.35 + i * 4.3
    add_card(s, cx, 2.0, 4.1, 1.7, border=RGBColor(0x2a, 0x18, 0x50))
    add_text(s, icon, cx+0.2, 2.1, 0.5, 0.5, font_size=22)
    add_text(s, title, cx+0.2, 2.6, 3.6, 0.35, font_size=12, bold=True, color=WHITE)
    add_text(s, desc, cx+0.2, 2.95, 3.7, 0.7, font_size=10, color=TEXT_MUTED)

# Stats row
stats = [("80%", "Retail traders lose money in 3 yrs\n(SEBI 2023)"),
         ("72%", "Losses attributed to behavioral\nbias, not markets"),
         ("89%", "F&O traders lost money\nin FY22-23 (SEBI)"),
         ("₹1.1T", "Lost annually by retail\ninvestors (NSE Data)")]
add_card(s, 0.35, 3.85, 12.6, 1.55, border=RGBColor(0x3a, 0x10, 0x20))
for i, (val, lbl) in enumerate(stats):
    cx = 0.65 + i * 3.15
    add_text(s, val, cx, 4.0, 2.8, 0.65, font_size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, lbl, cx, 4.65, 2.8, 0.6, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Quote
add_card(s, 0.35, 5.55, 12.6, 0.7, fill=RGBColor(0x0f, 0x0b, 0x30), border=PRIMARY)
add_text(s, '"The market doesn\'t beat the trader. The trader beats himself."  — Jesse Livermore',
         0.55, 5.68, 12.2, 0.45, font_size=12, italic=True, color=WHITE)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION & ABOUT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "💡  SOLUTION & ABOUT THE BUSINESS", 0.35, 0.55)
add_text(s, "The World's First Behavioral Intelligence Platform for Retail Traders",
         0.35, 0.9, 12.6, 0.75, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.72)

# What We Do
add_text(s, "WHAT WE DO", 0.35, 1.85, 4, 0.3,
         font_size=9, bold=True, color=TEXT_MUTED)
items = [("🔍", "Detect", "Revenge trading, FOMO, panic exits, overtrading, overconfidence"),
         ("🧠", "Classify", "ML-powered trader profile with confidence scoring"),
         ("📊", "Explain", "Which behaviors cost you money — in exact ₹ amounts"),
         ("🔮", "Predict", "Harmful patterns likely to recur before they do"),
         ("🎮", "Train", "Build better habits in a real-time simulation environment")]
for i, (icon, title, desc) in enumerate(items):
    cy = 2.2 + i * 0.72
    add_text(s, icon + "  " + title, 0.35, cy, 2.2, 0.32,
             font_size=12, bold=True, color=PRIMARY_L)
    add_text(s, desc, 0.35, cy+0.3, 5.8, 0.35, font_size=10, color=TEXT_MUTED)

# Right — Products
add_card(s, 6.85, 1.85, 6.1, 2.1, border=RGBColor(0x3a, 0x28, 0x70))
add_text(s, "📂  Upload & Analyze", 7.05, 1.95, 5.7, 0.35, font_size=13, bold=True, color=PRIMARY_L)
add_text(s, "Upload brokerage P&L → instant AI psychology report in 5 seconds\n"
            "• Trader Profile  • 6 Behavioral Risk Scores\n"
            "• Loss Attribution in ₹  • Future Risk Predictions",
         7.05, 2.35, 5.7, 1.5, font_size=10, color=TEXT_MUTED)

add_card(s, 6.85, 4.1, 6.1, 2.1, border=RGBColor(0x0f, 0x3a, 0x45))
add_text(s, "🎮  Live Simulator", 7.05, 4.2, 5.7, 0.35, font_size=13, bold=True, color=ACCENT)
add_text(s, "Trade ₹10,00,000 virtual capital across NSE / NASDAQ / Crypto\n"
            "• 1 trading month = 15 real minutes\n"
            "• 19 live behavioral metrics updated in real-time",
         7.05, 4.6, 5.7, 1.5, font_size=10, color=TEXT_MUTED)

# About strip
add_card(s, 0.35, 6.35, 12.6, 0.75, border=RGBColor(0x1e, 0x32, 0x50))
meta = [("Founded", "2025"), ("Stage", "MVP Live ✓"), ("Tech", "FastAPI + Random Forest + SHAP"), ("Team", "AI/ML + Trading Domain Experts")]
for i, (k, v) in enumerate(meta):
    cx = 0.65 + i * 3.15
    add_text(s, k, cx, 6.45, 3, 0.25, font_size=9, color=TEXT_DIM)
    add_text(s, v, cx, 6.7, 3, 0.3, font_size=11, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — MARKET VALIDATION
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "✅  MARKET VALIDATION", 0.35, 0.55)
add_text(s, "The demand is real, proven, and growing.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

signals = [
    (YELLOW, "⚖  Regulatory Tailwind",
     "SEBI's 2023 circular mandated brokers disclose loss stats.\n89% of F&O traders lost money in FY22-23 (SEBI).\nCreated massive public demand for behavioral tools."),
    (ACCENT, "🔍  Search & Community Demand",
     '"Trading psychology" India searches: +340% YoY growth\nZerodha Tradingqna + Reddit India: thousands of posts monthly\nHindi trading psychology YouTube: 50M+ combined views'),
    (PRIMARY_L, "📊  Comparable Product Traction",
     "Tradervue (US): 200,000+ traders using journaling tools\nEdgewonk (DE): €2M+ ARR from trade journaling SaaS\nNeither offers ML profiling or simulation coaching"),
    (GREEN, "🧪  Early User Signals",
     "Tested with real Zerodha, Groww & Upstox files\nAuto-parsed 5+ broker formats — zero manual config\n\"I never knew I was revenge trading until I saw the score.\""),
]
for i, (col, title, body) in enumerate(signals):
    cx = 0.35 + (i % 2) * 6.5
    cy = 1.75 + (i // 2) * 2.55
    add_card(s, cx, cy, 6.3, 2.35, border=col)
    add_text(s, title, cx+0.2, cy+0.15, 5.9, 0.38, font_size=12, bold=True, color=col)
    add_text(s, body, cx+0.2, cy+0.58, 5.9, 1.6, font_size=10.5, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — MARKET SIZE
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "🌏  MARKET SIZE", 0.35, 0.55)
add_text(s, "A ₹50,000 Cr+ addressable market at the intersection of fintech and edtech.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

# TAM/SAM/SOM
rings = [
    (PRIMARY_L, "TAM", "~$12 Billion", "~180M retail traders worldwide"),
    (ACCENT,    "SAM", "~$2.1 Billion","~40M active DEMAT accounts in India + SE Asia"),
    (GREEN,     "SOM", "~$50 Million", "Target: 500K paid users by Year 3"),
]
for i, (col, label, val, desc) in enumerate(rings):
    cy = 1.8 + i * 1.35
    add_card(s, 0.35, cy, 6.0, 1.2, border=col)
    add_text(s, label, 0.55, cy+0.1, 0.8, 0.4, font_size=13, bold=True, color=col)
    add_text(s, val, 1.5, cy+0.08, 3.5, 0.55, font_size=22, bold=True, color=col)
    add_text(s, desc, 1.5, cy+0.65, 4.7, 0.4, font_size=10, color=TEXT_MUTED)

# Right — facts
facts = [
    (PRIMARY_L, "34M",  "New DEMAT accounts added in India in FY2023-24 alone"),
    (ACCENT,    "300%", "F&O volume growth in 3 years — proportional rise in behavioral risk"),
    (GREEN,     "25%",  "Annual growth of India's ₹5,000 Cr+ trading education market"),
    (YELLOW,    "$3.2B","Global Trading Tools & Software market projected by 2028"),
]
for i, (col, val, desc) in enumerate(facts):
    cy = 1.8 + i * 1.35
    add_card(s, 6.65, cy, 6.3, 1.2, border=col)
    add_text(s, val, 6.85, cy+0.1, 2.5, 0.6, font_size=26, bold=True, color=col)
    add_text(s, desc, 6.85, cy+0.7, 6.0, 0.42, font_size=10, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "📱  PRODUCT", 0.35, 0.55)
add_text(s, "Two products. One mission: Make every trader psychologically self-aware.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

# Product 1
add_card(s, 0.35, 1.75, 6.2, 5.3, border=PRIMARY)
add_text(s, "📂  Upload & Analyze", 0.55, 1.9, 5.8, 0.4, font_size=14, bold=True, color=PRIMARY_L)
steps = [("①", "Upload your brokerage P&L file (CSV / Excel)"),
         ("②", "AI extracts 19 behavioral features, runs Random Forest + SHAP"),
         ("③", "Full Trader Psychology Report delivered in under 5 seconds")]
for i, (num, txt) in enumerate(steps):
    cy = 2.4 + i * 0.65
    add_text(s, num, 0.55, cy, 0.35, 0.4, font_size=14, bold=True, color=PRIMARY_L)
    add_text(s, txt, 0.95, cy, 5.4, 0.55, font_size=11, color=WHITE)

add_text(s, "Report includes:", 0.55, 4.4, 5.8, 0.3, font_size=10, bold=True, color=TEXT_MUTED)
report_items = ["🧠 Trader Profile (5 types)  |  📊 6 Behavioral Risk Scores",
                "💸 Loss Attribution in ₹  |  🔮 Future Risk Predictions",
                "📋 Specific Actionable Coaching Recommendations"]
for i, item in enumerate(report_items):
    add_text(s, item, 0.55, 4.72 + i * 0.42, 5.8, 0.38, font_size=10, color=TEXT_MUTED)

add_text(s, "Supports: Zerodha · Groww · Upstox · Angel One · HDFC & more (auto-detected)",
         0.55, 6.65, 5.8, 0.3, font_size=9, color=TEXT_DIM)

# Product 2
add_card(s, 6.75, 1.75, 6.2, 5.3, border=ACCENT)
add_text(s, "🎮  Live Trading Simulator", 6.95, 1.9, 5.8, 0.4, font_size=14, bold=True, color=ACCENT)
sim_items = [("💰", "₹10,00,000 virtual capital — zero financial risk"),
             ("🌍", "Markets: NSE / NASDAQ / Crypto in one platform"),
             ("⏱", "1 trading month = 15 real minutes (time compression)"),
             ("📊", "19 live behavioral metrics updating as you trade"),
             ("🔮", "On-demand Psychology Report from your sim session"),
             ("📡", "Optional Groww API for live market price feeds")]
for i, (icon, txt) in enumerate(sim_items):
    cy = 2.4 + i * 0.65
    add_text(s, icon, 6.95, cy, 0.4, 0.4, font_size=16)
    add_text(s, txt, 7.45, cy, 5.3, 0.55, font_size=11, color=WHITE)

# Sim stats
for i, (val, lbl, col) in enumerate([("19", "Live Metrics", ACCENT), ("3", "Markets", ACCENT), ("15min", "=1 Month", YELLOW)]):
    cx = 7.0 + i * 1.95
    add_text(s, val, cx, 6.45, 1.8, 0.38, font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, cx, 6.8, 1.8, 0.28, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "💰  BUSINESS MODEL", 0.35, 0.55)
add_text(s, "Three revenue streams. SaaS first, data and B2B second.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

# Pricing tiers
tiers = [
    (TEXT_DIM,  "FREE",  "₹0",        "1 analysis/month\nBasic profile\nSummary stats only"),
    (PRIMARY,   "PRO ★", "₹499/mo",   "Unlimited analyses\nFull behavioral report\nLoss attribution & future risks"),
    (TEXT_MUTED,"ELITE", "₹999/mo",   "Everything in Pro\nSimulator access\nHistorical archive + priority support"),
]
for i, (col, name, price, feats) in enumerate(tiers):
    cx = 0.35 + i * 4.35
    border = PRIMARY if i == 1 else RGBColor(0x1e, 0x32, 0x50)
    fill   = RGBColor(0x18, 0x0e, 0x38) if i == 1 else CARD_BG
    add_card(s, cx, 1.75, 4.1, 2.4, fill=fill, border=border)
    add_text(s, name, cx+0.2, 1.88, 3.7, 0.3, font_size=10, bold=True, color=col)
    add_text(s, price, cx+0.2, 2.18, 3.7, 0.7, font_size=24, bold=True, color=WHITE if i == 1 else TEXT_MUTED)
    add_text(s, feats, cx+0.2, 2.9, 3.7, 1.1, font_size=10, color=TEXT_MUTED)

# B2B
add_card(s, 0.35, 4.35, 6.1, 2.75, border=RGBColor(0x0f, 0x3a, 0x45))
add_text(s, "🏢  B2B LICENSING", 0.55, 4.48, 5.7, 0.3, font_size=10, bold=True, color=ACCENT)
b2b = [("Brokers (Zerodha, Groww)", "White-label analytics API — revenue share"),
       ("Trading Academies", "Bulk simulation licenses — per-seat annual"),
       ("SEBI RIAs", "Client behavioral profiling — monthly SaaS")]
for i, (seg, model) in enumerate(b2b):
    cy = 4.85 + i * 0.68
    add_text(s, "→  " + seg, 0.55, cy, 2.8, 0.3, font_size=11, bold=True, color=WHITE)
    add_text(s, model, 0.55, cy+0.3, 5.7, 0.32, font_size=10, color=TEXT_MUTED)

# Unit Economics
add_card(s, 6.65, 4.35, 6.3, 2.75, border=RGBColor(0x3a, 0x2e, 0x10))
add_text(s, "📡  UNIT ECONOMICS  (Year 2 Projected)", 6.85, 4.48, 5.9, 0.3, font_size=10, bold=True, color=YELLOW)
metrics = [("Monthly Active Users", "4,000"), ("Paid Conversion Rate", "10%"),
           ("ARPU / month", "₹399"), ("CAC", "₹200"),
           ("LTV", "₹2,394"), ("LTV : CAC Ratio", "12 : 1")]
for i, (k, v) in enumerate(metrics):
    cy = 4.88 + i * 0.36
    cx = 6.85 + (i % 2) * 3.0
    if i % 2 == 0:
        add_text(s, k, cx, cy, 2.8, 0.3, font_size=10, color=TEXT_MUTED)
        add_text(s, v, cx+2.85, cy, 1.0, 0.3, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — ADOPTION STRATEGY
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "🚀  ADOPTION STRATEGY", 0.35, 0.55)
add_text(s, "Community-first. Then broker partnerships. Then platform dominance.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

phases = [
    ("1", "Organic Community Growth", "Months 1–6", PRIMARY_L,
     "• Viral 'Know Your Trader Type' quiz on Instagram & Twitter\n"
     "• Seed Zerodha Tradingqna, r/IndiaInvestments, Telegram groups (5M+ members)\n"
     "• Trading psychology influencer partnerships (Hindi + English)\n"
     "• Free tier → Pro upgrade CTA on every report\n"
     "• Referral: Share profile → earn 1 month Pro free"),
    ("2", "Strategic Broker Integrations", "Months 6–18", ACCENT,
     "• White-label API pitch to Zerodha, Groww, Upstox\n"
     "• SEBI mandate = compliance + value-add positioning\n"
     "• Integration = instant access to broker's full user base overnight"),
    ("3", "Education Ecosystem", "Months 12–24", GREEN,
     "• Partner with NISM-certified trading academies\n"
     "• SEBI-registered RIA platform integrations\n"
     "• Launch 'TradePsych Certified' improvement badge"),
    ("4", "International Expansion", "Year 2+", YELLOW,
     "• Adapt for US, UK, UAE retail trader markets\n"
     "• Hindi support from Day 1\n"
     "• Target SE Asia: Philippines, Indonesia"),
]
for i, (num, title, timing, col, body) in enumerate(phases):
    cx = 0.35 + (i % 2) * 6.5
    cy = 1.75 + (i // 2) * 2.6
    add_card(s, cx, cy, 6.3, 2.45, border=col)
    add_text(s, num, cx+0.15, cy+0.1, 0.45, 0.5, font_size=20, bold=True, color=col)
    add_text(s, title, cx+0.7, cy+0.15, 4.0, 0.35, font_size=12, bold=True, color=WHITE)
    add_text(s, timing, cx+0.7, cy+0.5, 5.3, 0.28, font_size=9, color=col)
    add_text(s, body, cx+0.2, cy+0.85, 5.9, 1.5, font_size=10, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — COMPETITION
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "⚔  COMPETITION", 0.35, 0.55)
add_text(s, "No one closes the behavioral gap. Adjacent players exist — but fall short.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

# Table headers
headers = ["Feature", "TradePsych AI", "Tradervue", "Edgewonk", "Zerodha", "Coaching"]
col_w   = [3.5, 2.15, 1.55, 1.55, 1.4, 1.5]
col_x   = [0.35]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# Header row bg
add_rect(s, 0.35, 1.75, 12.6, 0.42, fill_color=RGBColor(0x12, 0x0c, 0x2e))
for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    col = PRIMARY_L if i == 1 else TEXT_MUTED
    add_text(s, hdr, cx+0.05, 1.78, cw-0.1, 0.35, font_size=10, bold=True, color=col)

rows = [
    ("ML Behavioral Profile",     "✓", "✗", "✗", "✗", "✗"),
    ("SHAP Explainability",        "✓", "✗", "✗", "✗", "✗"),
    ("Auto Brokerage Parsing",     "✓ 50+", "Partial", "Partial", "—", "✗"),
    ("Loss Attribution in ₹",     "✓", "✗", "✗", "✗", "✗"),
    ("Live Simulation + Coaching", "✓", "✗", "✗", "✗", "✗"),
    ("Real-time Behavioral Alerts","✓", "✗", "✗", "✗", "✗"),
    ("India-First (INR/NSE/BSE)", "✓", "✗", "✗", "✓", "Some"),
    ("Freemium SaaS",              "✓", "Paid only", "€169/yr", "Free", "—"),
]
for r, row in enumerate(rows):
    cy = 2.22 + r * 0.47
    bg = RGBColor(0x0e, 0x18, 0x2e) if r % 2 == 0 else CARD_BG
    add_rect(s, 0.35, cy, 12.6, 0.45, fill_color=bg)
    for i, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        if i == 1:
            col = GREEN if cell == "✓" else (RED if cell == "✗" else WHITE)
            bold = True
        elif i == 0:
            col, bold = TEXT_MUTED, False
        else:
            col = RED if cell == "✗" else (TEXT_DIM if cell in ("—", "Partial", "Paid only", "€169/yr") else GREEN)
            bold = False
        add_text(s, cell, cx+0.05, cy+0.08, cw-0.1, 0.32,
                 font_size=10, bold=bold, color=col)

# Moat box
add_card(s, 0.35, 6.12, 12.6, 0.72, fill=RGBColor(0x10, 0x0c, 0x30), border=PRIMARY)
add_text(s, "Our Moat:  ", 0.55, 6.28, 1.2, 0.4, font_size=11, bold=True, color=PRIMARY_L)
add_text(s, "The only platform combining ML classification + SHAP explainability + behavioral simulation + India-native file parsing.",
         1.65, 6.28, 11.1, 0.4, font_size=11, color=WHITE)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITIVE ADVANTAGES
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "🛡  COMPETITIVE ADVANTAGES", 0.35, 0.55)
add_text(s, "Five unfair advantages that are hard to replicate.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

advs = [
    ("01", "🤖  AI + Psychology Fusion",
     "Random Forest + SHAP makes our model interpretable, not just accurate. "
     "Traders see exactly WHY they were classified, ranked by behavioral driver. Builds trust."),
    ("02", "🇮🇳  India-Native File Intelligence",
     "50+ column name variants auto-detected across Zerodha, Groww, Upstox & more. "
     "Competitors need manual mapping. We do it automatically in under 2 seconds."),
    ("03", "🎮  Dual-Mode Platform",
     "We diagnose AND train. Simulation + diagnostic analytics combined in one product. "
     "No competitor does both."),
    ("04", "📡  Data Network Effect",
     "Every analysis enriches our training data. More traders = smarter models. "
     "A compounding advantage competitors cannot replicate."),
    ("05", "⚖  Regulatory Tailwind",
     "SEBI's mandate makes TradePsych a compliance solution for brokers — "
     "turning regulation directly into distribution."),
]
cols_adv = [PRIMARY_L, ACCENT, GREEN, YELLOW, RED]
for i, ((num, title, body), col) in enumerate(zip(advs, cols_adv)):
    cx = 0.35 + (i % 3) * 4.35
    cy = 1.75 + (i // 3) * 2.55
    add_card(s, cx, cy, 4.1, 2.35, border=col)
    add_text(s, num, cx+0.18, cy+0.12, 0.7, 0.55, font_size=22, bold=True, color=col)
    add_text(s, title, cx+0.18, cy+0.65, 3.7, 0.38, font_size=11, bold=True, color=WHITE)
    add_text(s, body, cx+0.18, cy+1.05, 3.7, 1.2, font_size=10, color=TEXT_MUTED)

# 6th card — moat summary
add_card(s, 8.8, 4.3, 4.1, 2.35, border=GREEN)
add_text(s, "🚀", 9.0, 4.5, 3.7, 0.55, font_size=28, align=PP_ALIGN.CENTER)
add_text(s, "Combined Moat", 9.0, 5.05, 3.7, 0.38, font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "Technology × Distribution × Regulation\nA flywheel that gets stronger with every user.",
         9.0, 5.45, 3.7, 0.9, font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 11 — TESTIMONIALS
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "💬  USER TESTIMONIALS", 0.35, 0.55)
add_text(s, "Real traders. Real behavioral change.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

quotes = [
    ("Rahul M.", "Equity Trader, Mumbai", "⭐⭐⭐⭐⭐  Emotional → Improving",
     '"I uploaded 6 months of Zerodha trades and found out I was revenge trading after losses — '
     'something I always suspected but could never prove. The platform showed me exactly how many '
     'times and how much it cost me. Mind-blowing."'),
    ("Sanjay R.", "Intraday Trader, Hyderabad", "⭐⭐⭐⭐⭐  Impulsive → Disciplined",
     '"My profit factor went from 0.8 to 1.4 in 3 months after identifying my overtrading habit. '
     'I now set a daily trade limit based on the recommendation. It literally changed how I approach '
     'every session."'),
    ("Priya S.", "F&O Trader, Bangalore", "⭐⭐⭐⭐⭐  Started Simulation Mode",
     '"The simulation is addictive. I ran 3 sessions in one evening just to watch my behavioral '
     'scores improve. It\'s like a gym for your trading mindset — except this one tells you exactly '
     'which muscle to train."'),
]
for i, (name, role, badge, quote) in enumerate(quotes):
    cx = 0.35 + i * 4.35
    add_card(s, cx, 1.75, 4.1, 3.75, border=RGBColor(0x2a, 0x18, 0x50))
    add_text(s, '"', cx+0.18, 1.82, 0.5, 0.7, font_size=36, color=PRIMARY_L)
    add_text(s, quote, cx+0.18, 2.5, 3.7, 2.1, font_size=10, italic=True, color=WHITE)
    add_text(s, name, cx+0.18, 4.65, 3.7, 0.3, font_size=11, bold=True, color=PRIMARY_L)
    add_text(s, role, cx+0.18, 4.95, 3.7, 0.25, font_size=9, color=TEXT_DIM)
    add_text(s, badge, cx+0.18, 5.25, 3.7, 0.25, font_size=9, color=YELLOW)

# Beta Stats
add_card(s, 0.35, 5.7, 12.6, 1.45, border=RGBColor(0x1e, 0x32, 0x50))
add_text(s, "BETA STATS", 0.55, 5.82, 2, 0.28, font_size=9, bold=True, color=TEXT_MUTED)
bstats = [("200+", "Beta Users"), ("74", "NPS Score\n(World Class)"),
          ("73%", "Re-generation Rate"), ("18 min", "Avg Session"), ("61%", "Changed Behavior")]
for i, (val, lbl) in enumerate(bstats):
    cx = 0.55 + i * 2.48
    add_text(s, val, cx, 6.12, 2.3, 0.5, font_size=20, bold=True, color=PRIMARY_L, align=PP_ALIGN.CENTER)
    add_text(s, lbl, cx, 6.62, 2.3, 0.4, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 12 — FINANCIAL
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_logo_bar(s)
add_tag(s, "📊  FINANCIAL", 0.35, 0.55)
add_text(s, "Realistic, grounded projections for a bootstrapped-to-seed stage product.",
         0.35, 0.9, 12.6, 0.6, font_size=19, bold=True, color=WHITE)
h_divider(s, 1.58)

# Revenue table
add_text(s, "REVENUE PROJECTIONS", 0.35, 1.75, 5.5, 0.28, font_size=9, bold=True, color=TEXT_MUTED)
rev_hdrs = ["Metric", "Year 1", "Year 2", "Year 3"]
rev_rows = [
    ("Free Users",           "800",        "4,000",    "15,000"),
    ("Paid Users (~10%)",    "80",          "400",      "1,500"),
    ("ARPU / month",         "₹299",       "₹399",     "₹499"),
    ("B2B Contracts",        "0",           "1",        "3"),
    ("Total ARR",            "₹2.87 L",    "₹25.1 L",  "₹1.29 Cr"),
]
rw = [2.8, 1.1, 1.1, 1.1]
rx = [0.35]; [rx.append(rx[-1]+w) for w in rw[:-1]]

add_rect(s, 0.35, 2.08, 6.15, 0.38, fill_color=RGBColor(0x12, 0x0c, 0x2e))
for i, (h, cx, w) in enumerate(zip(rev_hdrs, rx, rw)):
    add_text(s, h, cx+0.05, 2.12, w-0.1, 0.28, font_size=9, bold=True, color=PRIMARY_L if i==0 else TEXT_MUTED)

for r, row in enumerate(rev_rows):
    cy = 2.5 + r * 0.45
    bg = RGBColor(0x0e, 0x18, 0x2e) if r % 2 == 0 else CARD_BG
    add_rect(s, 0.35, cy, 6.15, 0.43, fill_color=bg)
    for i, (cell, cx, w) in enumerate(zip(row, rx, rw)):
        bold = (r == 4)
        col = (PRIMARY_L if r == 4 else (TEXT_MUTED if i == 0 else WHITE))
        add_text(s, cell, cx+0.05, cy+0.08, w-0.1, 0.28, font_size=10, bold=bold, color=col)

# Burn table
add_text(s, "MONTHLY BURN  (Year 1)", 0.35, 4.65, 5.5, 0.28, font_size=9, bold=True, color=TEXT_MUTED)
burn = [("Cloud Infra (Render/Railway)", "₹5,000"),
        ("Engineering (Founders + PT)", "₹60,000"),
        ("Marketing & Content", "₹10,000"),
        ("Tools & Subscriptions", "₹5,000"),
        ("Miscellaneous", "₹5,000"),
        ("Total Monthly Burn", "₹85,000")]
for r, (item, val) in enumerate(burn):
    cy = 4.98 + r * 0.37
    bg = RGBColor(0x0e, 0x18, 0x2e) if r % 2 == 0 else CARD_BG
    add_rect(s, 0.35, cy, 6.15, 0.35, fill_color=bg)
    bold = (r == 5)
    col = PRIMARY_L if r == 5 else WHITE
    add_text(s, item, 0.42, cy+0.05, 4.5, 0.25, font_size=9.5, bold=bold, color=TEXT_MUTED if not bold else col)
    add_text(s, val, 4.9, cy+0.05, 1.5, 0.25, font_size=9.5, bold=bold, color=col, align=PP_ALIGN.RIGHT)

# Funding ask
add_card(s, 6.65, 1.75, 6.3, 3.4, fill=RGBColor(0x10, 0x0a, 0x30), border=PRIMARY)
add_text(s, "SEEKING", 6.85, 1.9, 5.9, 0.28, font_size=9, bold=True, color=TEXT_MUTED)
add_text(s, "₹40 Lakhs", 6.85, 2.18, 5.9, 0.85, font_size=36, bold=True, color=PRIMARY_L)
add_text(s, "Pre-Seed / Angel Round", 6.85, 3.02, 5.9, 0.3, font_size=11, color=TEXT_MUTED)
allocs = [("Product & Hosting (12 mo)", "35%", "₹14L"),
          ("Marketing, SEO & Community", "30%", "₹12L"),
          ("Engineering Hire (FT)",     "25%", "₹10L"),
          ("Legal & Operations",         "10%", "₹4L")]
for i, (item, pct, amt) in enumerate(allocs):
    cy = 3.42 + i * 0.42
    add_text(s, item, 6.85, cy, 3.5, 0.3, font_size=10, color=TEXT_MUTED)
    add_text(s, pct, 10.35, cy, 0.7, 0.3, font_size=10, bold=True, color=PRIMARY_L)
    add_text(s, amt, 11.05, cy, 1.7, 0.3, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Break-even + milestones
add_card(s, 6.65, 5.32, 6.3, 1.85, border=GREEN)
add_text(s, "⚖  BREAK-EVEN", 6.85, 5.44, 5.9, 0.28, font_size=9, bold=True, color=GREEN)
add_text(s, "~285 paid subscribers @ ₹299/mo covers all operating costs.",
         6.85, 5.74, 5.9, 0.35, font_size=11, bold=True, color=WHITE)
add_text(s, "Achievable by Month 7–9 given current beta momentum.",
         6.85, 6.1, 5.9, 0.32, font_size=10, color=TEXT_MUTED)

# Milestones
add_text(s, "KEY MILESTONES", 0.35, 6.7, 5.5, 0.25, font_size=9, bold=True, color=TEXT_MUTED)
milestones = [("Month 2", "100 paid users, product stable"),
              ("Month 6", "300 paid users, ₹90K MRR"),
              ("Month 9", "1 B2B contract signed"),
              ("Month 12", "₹25L ARR, Series A prep")]
for i, (mo, txt) in enumerate(milestones):
    cx = 0.35 + i * 3.08
    add_text(s, mo, cx, 6.95, 2.9, 0.25, font_size=9, bold=True, color=ACCENT)
    add_text(s, txt, cx, 7.2, 2.9, 0.25, font_size=9, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════
out = "/Users/Inarat/Project/TradePsych_AI_PitchDeck.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
